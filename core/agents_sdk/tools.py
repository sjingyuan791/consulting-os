"""
OS 既存関数を @function_tool として定義するモジュール。

ツール割当制限（不変条件）:
  CFO Agent    : get_financials_verified, python_reconcile_financials
  Writer Agent : check_section_schema, generate_docx, generate_pptx, generate_pdf_report
  Skeptic Agent: run_quality_gate, run_fact_check, create_human_checkpoint
  Strategy Agent: （ツールなし — LLM 推論のみ）

エラー時は必ずフェイルセーフ値を返す（例外を伝播させない）。
"""
from __future__ import annotations

import json
import logging
import os
from datetime import datetime
from typing import Optional

from agents import function_tool

logger = logging.getLogger("agents_sdk.tools")

BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
ARTIFACTS_DIR = os.path.join(BASE_DIR, "artifacts_out")


# ============================================================
# CFO Agent Tools（制限: この2つのみ）
# ============================================================

@function_tool
def get_financials_verified(pipeline_run_id: str) -> str:
    """
    Check whether financial data has been verified for the given pipeline run.
    MUST be the first tool called by the CFO Agent.
    If financials_verified is False, the CFO Agent must stop and return
    missing_inputs=['financials_verified'] immediately.

    Args:
        pipeline_run_id: UUID of the pipeline run to check.

    Returns:
        JSON string: {"financials_verified": bool, "pipeline_run_id": str}
    """
    try:
        from core.supabase_client import get_supabase_client
        sb = get_supabase_client()

        # pipeline_runs → client_id を取得
        pr_res = (
            sb.table("pipeline_runs")
            .select("client_id")
            .eq("id", pipeline_run_id)
            .limit(1)
            .execute()
        )
        if not pr_res.data:
            return json.dumps({"financials_verified": False,
                               "pipeline_run_id": pipeline_run_id,
                               "reason": "pipeline_run not found"})

        client_id = pr_res.data[0]["client_id"]

        # analysis_runs の最新レコードから financials_verified を取得
        ar_res = (
            sb.table("analysis_runs")
            .select("financials_verified")
            .eq("client_id", client_id)
            .order("created_at", desc=True)
            .limit(1)
            .execute()
        )
        if ar_res.data:
            verified = bool(ar_res.data[0].get("financials_verified", False))
        else:
            verified = False

        return json.dumps({"financials_verified": verified,
                           "pipeline_run_id": pipeline_run_id})
    except Exception as exc:
        logger.error("get_financials_verified failed: %s", exc)
        # フェイルセーフ: 未検証として返す
        return json.dumps({"financials_verified": False,
                           "pipeline_run_id": pipeline_run_id,
                           "error": str(exc)})


@function_tool
def python_reconcile_financials(financials_json: str) -> str:
    """
    Run deterministic reconciliation checks on financial data.
    Only call this after get_financials_verified returns True.

    Checks:
    - BS balance: total_assets == total_liabilities + net_assets (within 1%)
    - Gross profit: sales - cogs == gross_profit (within 1%)

    Args:
        financials_json: JSON string with keys: total_assets, total_liabilities,
            net_assets, sales, cogs, gross_profit (all numeric, JPY).

    Returns:
        JSON string: {"passed": bool, "items": [...], "errors": [...]}
    """
    try:
        data = json.loads(financials_json)
    except json.JSONDecodeError as exc:
        return json.dumps({"passed": False, "items": [],
                           "errors": [f"JSON parse error: {exc}"]})

    items: list = []
    errors: list = []

    def _get(key: str) -> Optional[float]:
        v = data.get(key)
        return float(v) if v is not None else None

    # BS バランスチェック
    assets = _get("total_assets")
    liab = _get("total_liabilities")
    equity = _get("net_assets")
    if all(v is not None for v in [assets, liab, equity]):
        expected = liab + equity  # type: ignore[operator]
        delta = abs(assets - expected)  # type: ignore[operator]
        tol = (assets or 1) * 0.01
        passed = delta <= tol
        items.append({"check": "BS_balance",
                      "result": "pass" if passed else "fail",
                      "delta": round(delta, 0)})
        if not passed:
            errors.append(
                f"BS不一致: 総資産({assets:,.0f}) ≠ 負債+純資産({expected:,.0f})"
            )

    # 粗利整合チェック
    sales = _get("sales")
    cogs = _get("cogs")
    gross = _get("gross_profit")
    if all(v is not None for v in [sales, cogs, gross]):
        expected_gross = sales - cogs  # type: ignore[operator]
        delta = abs(gross - expected_gross)  # type: ignore[operator]
        tol = (sales or 1) * 0.01
        passed = delta <= tol
        items.append({"check": "gross_profit_consistency",
                      "result": "pass" if passed else "fail",
                      "delta": round(delta, 0)})
        if not passed:
            errors.append(
                f"粗利不一致: {gross:,.0f} ≠ 売上-原価({expected_gross:,.0f})"
            )

    return json.dumps({"passed": len(errors) == 0, "items": items, "errors": errors},
                      ensure_ascii=False)


# ============================================================
# Writer Agent Tools
# ============================================================

@function_tool
def check_section_schema(section_id: int, data_json: str) -> str:
    """
    Validate that a section's data dict conforms to the expected Pydantic schema
    for section_id (1-13).

    Args:
        section_id: Section number (1-13).
        data_json: JSON string of the section data to validate.

    Returns:
        JSON string: {"valid": bool, "schema_name": str, "errors": [...]}
    """
    try:
        from core.midterm_plan_engine import SECTION_ID_TO_MODEL
        model_class = SECTION_ID_TO_MODEL.get(section_id)
        if not model_class:
            return json.dumps({"valid": False, "schema_name": "unknown",
                               "errors": [f"No schema for section_id={section_id}"]})

        try:
            data = json.loads(data_json)
        except json.JSONDecodeError as exc:
            return json.dumps({"valid": False,
                               "schema_name": model_class.__name__,
                               "errors": [f"JSON parse error: {exc}"]})

        try:
            model_class(**data)
            return json.dumps({"valid": True,
                               "schema_name": model_class.__name__,
                               "errors": []})
        except Exception as exc:
            from pydantic import ValidationError
            if isinstance(exc, ValidationError):
                errs = [f"{e['loc']}: {e['msg']}" for e in exc.errors()]
            else:
                errs = [str(exc)]
            return json.dumps({"valid": False,
                               "schema_name": model_class.__name__,
                               "errors": errs})
    except Exception as exc:
        logger.error("check_section_schema failed: %s", exc)
        return json.dumps({"valid": False, "schema_name": "error",
                           "errors": [str(exc)]})


@function_tool
def generate_docx(sections_json: str, client_name: str, output_filename: str) -> str:
    """
    Generate a DOCX file from mid-term plan sections using DocxWriter.

    Args:
        sections_json: JSON array of {"title": str, "narrative": str} objects.
        client_name: Company name for the document title.
        output_filename: Filename without extension (saved to artifacts_out/).

    Returns:
        JSON string: {"success": bool, "output_path": str, "error": str}
    """
    try:
        sections = json.loads(sections_json)
    except json.JSONDecodeError as exc:
        return json.dumps({"success": False, "output_path": "",
                           "error": f"JSON parse: {exc}"})
    try:
        from core.docx_writer import DocxWriter
        writer = DocxWriter()
        writer.add_title(f"中期経営計画書 — {client_name}")
        for sec in sections:
            writer.add_chapter(sec.get("title", ""))
            writer.add_content_from_markdown(sec.get("narrative", ""))

        os.makedirs(ARTIFACTS_DIR, exist_ok=True)
        path = os.path.join(ARTIFACTS_DIR, f"{output_filename}.docx")
        writer.save(path)
        return json.dumps({"success": True, "output_path": path, "error": ""})
    except Exception as exc:
        logger.error("generate_docx failed: %s", exc)
        return json.dumps({"success": False, "output_path": "", "error": str(exc)})


@function_tool
def generate_pptx(sections_json: str, client_name: str, output_filename: str) -> str:
    """
    Generate a PowerPoint file from mid-term plan sections using PPTWriter.

    Args:
        sections_json: JSON array of {"title": str, "bullets": [str]} objects.
        client_name: Company name for the title slide.
        output_filename: Filename without extension (saved to artifacts_out/).

    Returns:
        JSON string: {"success": bool, "output_path": str, "error": str}
    """
    try:
        sections = json.loads(sections_json)
    except json.JSONDecodeError as exc:
        return json.dumps({"success": False, "output_path": "",
                           "error": f"JSON parse: {exc}"})
    try:
        from pptx import Presentation
        prs = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "中期経営計画書"
        title_slide.placeholders[1].text = client_name

        for sec in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec.get("title", "")
            tf = slide.placeholders[1].text_frame
            tf.text = ""
            for bullet in sec.get("bullets", []):
                p = tf.add_paragraph()
                p.text = str(bullet)

        os.makedirs(ARTIFACTS_DIR, exist_ok=True)
        path = os.path.join(ARTIFACTS_DIR, f"{output_filename}.pptx")
        prs.save(path)
        return json.dumps({"success": True, "output_path": path, "error": ""})
    except Exception as exc:
        logger.error("generate_pptx failed: %s", exc)
        return json.dumps({"success": False, "output_path": "", "error": str(exc)})


@function_tool
def generate_pdf_report(report_data_json: str, output_filename: str) -> str:
    """
    Generate a PDF report using the existing ReportLab writer.

    Args:
        report_data_json: JSON string matching core/report_pdf.py generate_pdf_report() format.
        output_filename: Filename without extension (saved to artifacts_out/).

    Returns:
        JSON string: {"success": bool, "output_path": str, "error": str}
    """
    try:
        data = json.loads(report_data_json)
    except json.JSONDecodeError as exc:
        return json.dumps({"success": False, "output_path": "",
                           "error": f"JSON parse: {exc}"})
    try:
        from core.report_pdf import generate_pdf_report as _gen_pdf
        os.makedirs(ARTIFACTS_DIR, exist_ok=True)
        path = os.path.join(ARTIFACTS_DIR, f"{output_filename}.pdf")
        _gen_pdf(data, path)
        return json.dumps({"success": True, "output_path": path, "error": ""})
    except Exception as exc:
        logger.error("generate_pdf_report failed: %s", exc)
        return json.dumps({"success": False, "output_path": "", "error": str(exc)})


# ============================================================
# Skeptic Agent Tools
# ============================================================

@function_tool
def run_quality_gate(plan_json: str) -> str:
    """
    Run the Decision-Grade quality gate on a RefinedStrategicPlan.
    Wraps check_strategic_refinement_quality() from quality_gate_enhanced.py.

    Args:
        plan_json: JSON string of the RefinedStrategicPlan model dict.

    Returns:
        JSON string: {"status": "approved"|"blocked"|"warning",
                      "blocking_reasons": [...], "warnings": [...]}
    """
    try:
        from core.schemas.refinement_schema import RefinedStrategicPlan
        from core.quality_gate_enhanced import check_strategic_refinement_quality

        plan_data = json.loads(plan_json)
        plan = RefinedStrategicPlan(**plan_data)
        result = check_strategic_refinement_quality(plan)
        return json.dumps({
            "status": result.status,
            "blocking_reasons": result.blocking_reasons,
            "warnings": result.warnings,
        }, ensure_ascii=False)
    except Exception as exc:
        logger.error("run_quality_gate failed: %s", exc)
        return json.dumps({
            "status": "blocked",
            "blocking_reasons": [f"Quality gate execution error: {exc}"],
            "warnings": [],
        }, ensure_ascii=False)


@function_tool
def run_fact_check(ai_response: str, source_data_json: str) -> str:
    """
    Run FactCheckEngine against an AI-generated response.
    Use this to verify that citations in the response match source data.

    AI output must use 【】citation markers (e.g., 【財務データ:売上高2024年】)
    for the fact-checker to find and verify them.

    Args:
        ai_response: The AI-generated text to fact-check.
        source_data_json: JSON string of the source data dict.

    Returns:
        JSON string: {"trust_score": float, "total_citations": int,
                      "verified_count": int, "mismatch_count": int,
                      "hallucination_count": int, "requires_human_review": bool,
                      "review_points": [...], "critical_issues": [...]}
    """
    try:
        source_data = json.loads(source_data_json)
    except json.JSONDecodeError as exc:
        return json.dumps({
            "trust_score": 0, "total_citations": 0, "verified_count": 0,
            "mismatch_count": 0, "hallucination_count": 0,
            "requires_human_review": True,
            "critical_issues": [f"Source data JSON error: {exc}"],
            "review_points": [],
        })
    try:
        from core.ai_quality_assurance import FactCheckEngine
        engine = FactCheckEngine(source_data)
        report = engine.check(ai_response, response_id="skeptic_check")
        return json.dumps({
            "trust_score": report.trust_score,
            "total_citations": report.total_citations,
            "verified_count": report.verified_count,
            "mismatch_count": report.mismatch_count,
            "hallucination_count": report.hallucination_count,
            "requires_human_review": report.requires_human_review,
            "review_points": report.review_points,
            "critical_issues": report.critical_issues,
        }, ensure_ascii=False)
    except Exception as exc:
        logger.error("run_fact_check failed: %s", exc)
        return json.dumps({
            "trust_score": 0, "total_citations": 0, "verified_count": 0,
            "mismatch_count": 0, "hallucination_count": 0,
            "requires_human_review": True,
            "critical_issues": [f"FactCheck execution error: {exc}"],
            "review_points": [],
        })


@function_tool
def create_human_checkpoint(
    pipeline_run_id: str,
    stage_output_id: str,
    blocking_reasons_json: str,
) -> str:
    """
    Create a human_checkpoints row when the Skeptic Agent blocks a plan.
    The pipeline will be paused until a human reviews and approves.

    Args:
        pipeline_run_id: UUID of the pipeline run.
        stage_output_id: UUID of the stage_output that triggered the block.
        blocking_reasons_json: JSON array string of reason strings.

    Returns:
        JSON string: {"checkpoint_id": str, "status": str, "error": str}
    """
    try:
        reasons = json.loads(blocking_reasons_json) if blocking_reasons_json else []
    except json.JSONDecodeError:
        reasons = [blocking_reasons_json]

    try:
        from core.supabase_client import get_supabase_client
        sb = get_supabase_client()
        payload = {
            "stage_output_id": stage_output_id,
            "pipeline_run_id": pipeline_run_id,
            "checkpoint_type": "skeptic_block",
            "checkpoint_name": "Skeptic Agent Decision-Grade Block",
            "status": "pending",
            "feedback_json": {"blocking_reasons": reasons},
            "notification_sent_at": datetime.now().isoformat(),
        }
        res = sb.table("human_checkpoints").insert(payload).execute()
        checkpoint_id = res.data[0]["id"] if res.data else ""
        return json.dumps({"checkpoint_id": checkpoint_id,
                           "status": "pending", "error": ""})
    except Exception as exc:
        logger.error("create_human_checkpoint failed: %s", exc)
        return json.dumps({"checkpoint_id": "", "status": "error",
                           "error": str(exc)})
