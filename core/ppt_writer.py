from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from core.models import StrategyContext
from core.schemas.midterm_plan_schema import MidtermPlanDocument

class PPTWriter:
    def __init__(self):
        self.prs = Presentation()

    def create_slides(self, context: StrategyContext, package_data: dict):
        # 1. Title Slide
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Strategic Revitalization Plan"
        subtitle.text = f"Prepared for {context.company_summary}"

        # Extract data from package dict safely
        diag = package_data.get("root_cause_diagnosis", {})
        opts = package_data.get("strategy_options", {})
        roadmap = package_data.get("execution_roadmap", {})
        
        # 2. Executive Summary (from Diagnosis Root Issue)
        root_issue = diag.get("root_issue", "N/A")
        primary = diag.get("primary_symptom", "")
        self._add_bullet_slide("Executive Summary (Diagnosis)", [f"Core Issue: {root_issue}", f"Primary Symptom: {primary}"])

        # 3. Key Issues (Hypotheses / Likely Causes)
        causes = diag.get("likely_root_causes", [])
        self._add_bullet_slide("Key Strategic Issues (Root Causes)", causes)

        # 4. Strategy Options & Selection
        selected = package_data.get("selected_strategy", {})
        chosen_id = selected.get("chosen_option_id")
        
        # Find the chosen option details
        chosen_opt = next((o for o in opts.get("options", []) if o.get("id") == chosen_id), None)
        
        if chosen_opt:
            self._add_bullet_slide(f"Selected Strategy: {chosen_opt.get('name')}", [
                f"Description: {chosen_opt.get('description')}",
                f"Rationale: {chosen_opt.get('rationale')}"
            ])
            
        # 5. Roadmap / Actions
        phases = roadmap.get("timeline", [])
        action_summary = []
        for p in phases:
             action_summary.append(f"{p.get('name')} ({p.get('duration')}): {', '.join(p.get('key_actions', []))}")
             
        self._add_bullet_slide("Execution Roadmap", action_summary)
        
        # 5. Risks & Opportunities
        self._add_bullet_slide("Risks", context.risks)
        self._add_bullet_slide("Opportunities", context.opportunities)

    def create_slides_from_midterm_plan(self, document: MidtermPlanDocument):
        """13セクションの中期経営計画書からスライドを生成"""
        # 1. Title Slide
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Mid-Term Management Plan"
        subtitle.text = f"Prepared for {document.client_id}\nPeriod: {document.plan_period}"

        # 2. Add slides for each section
        for section in document.sections:
            # Skip empty sections if any
            if not section.narrative:
                continue
                
            # Naive approach: Create one slide per section title, with truncated narrative items
            # Parse markdown roughly to get headers or bullets
            # For MVP, just put the first few paragraphs or bullet points
            
            content_lines = []
            for line in section.narrative.split("\n"):
                clean = line.strip()
                if clean.startswith(("- ", "* ")):
                    content_lines.append(clean[2:])
                elif clean.startswith("###"):
                    content_lines.append(clean.replace("#", "").strip())
                elif clean and not clean.startswith("#"):
                    # Limit paragraph length
                    content_lines.append(clean[:100] + "..." if len(clean) > 100 else clean)
                
                if len(content_lines) > 6: # limit per slide
                    break
            
            self._add_bullet_slide(f"{section.section_id}. {section.section_title}", content_lines)

    def _add_bullet_slide(self, title_text, content_items):
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        if isinstance(content_items, str):
            content_items = [content_items]
            
        for item in content_items:
            p = tf.add_paragraph()
            p.text = str(item)
            p.level = 0

    def save(self, path):
        self.prs.save(path)
