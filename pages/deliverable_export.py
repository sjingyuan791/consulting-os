"""
deliverable_export.py — 成果物エクスポートページ (STEP 用ドキュメント生成)

ProjectContext の全成果物を PPTX / DOCX に自動組み立てしてダウンロード提供する。
"""
import io
import tempfile
import datetime
import streamlit as st

from core.auth import check_auth
from core.sidebar import render_sidebar
from core.project_context import ProjectContext

# ------------------------------------------------------------------ #
#  python-pptx helpers
# ------------------------------------------------------------------ #

def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ConsultingPPT:
    """コンサルPPTX ビルダー"""

    NAVY   = "#1e3a5f"
    WHITE  = "#ffffff"
    LIGHT  = "#f0f4f8"
    ACCENT = "#2563eb"
    TEXT   = "#1e293b"

    def __init__(self, company_name: str):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        self.prs = Presentation()
        self.prs.slide_width  = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.company_name = company_name
        self._blank_layout = self.prs.slide_layouts[6]  # blank

    # ---- low-level drawing helpers ---------------------------------- #

    def _rect(self, slide, left, top, width, height, fill_hex, line_hex=None):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(fill_hex)
        if line_hex:
            shape.line.color.rgb = _rgb(line_hex)
        else:
            shape.line.fill.background()
        return shape

    def _textbox(self, slide, left, top, width, height, text,
                 font_size=14, bold=False, color_hex=None, wrap=True, align="left"):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color_hex:
            run.font.color.rgb = _rgb(color_hex)
        return txb

    def _header_bar(self, slide, title: str, subtitle: str = ""):
        from pptx.util import Pt
        # dark header band
        self._rect(slide, 0, 0, 13.33, 1.2, self.NAVY)
        self._textbox(slide, 0.3, 0.1, 11, 0.6, title,
                      font_size=24, bold=True, color_hex=self.WHITE)
        if subtitle:
            self._textbox(slide, 0.3, 0.75, 11, 0.4, subtitle,
                          font_size=13, color_hex="#90b8e0")
        # company name right
        self._textbox(slide, 10, 0.1, 3.0, 0.5, self.company_name,
                      font_size=11, color_hex="#90b8e0", align="center")

    def _bullets(self, slide, items: list, left=0.5, top=1.4, width=12.3, font_size=13):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(7.5 - top - 0.3))
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = f"・ {item}"
            p.space_before = Pt(4)
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = _rgb(self.TEXT)

    def _table_slide(self, slide, headers: list, rows: list, top=1.4, font_size=11):
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        cols = len(headers)
        if not rows:
            return
        tbl = slide.shapes.add_table(
            len(rows) + 1, cols,
            Inches(0.4), Inches(top),
            Inches(12.5), Inches(min(5.7, 0.35 * (len(rows) + 1)))
        ).table
        # header row
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(self.NAVY)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # data rows
        for ri, row in enumerate(rows):
            bg = self.LIGHT if ri % 2 == 0 else self.WHITE
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val or "")
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(bg)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = _rgb(self.TEXT)

    # ---- slide builders --------------------------------------------- #

    def add_title_slide(self, date_str: str):
        slide = self.prs.slides.add_slide(self._blank_layout)
        # full background
        self._rect(slide, 0, 0, 13.33, 7.5, self.NAVY)
        self._rect(slide, 0, 5.2, 13.33, 2.3, self.ACCENT)
        self._textbox(slide, 1, 1.5, 11.3, 1.2,
                      "中期経営計画 策定支援レポート",
                      font_size=32, bold=True, color_hex=self.WHITE, align="center")
        self._textbox(slide, 1, 2.8, 11.3, 0.8,
                      self.company_name,
                      font_size=22, color_hex="#90b8e0", align="center")
        self._textbox(slide, 1, 5.5, 11.3, 0.6,
                      f"作成日: {date_str}　　Consulting OS 自動生成レポート",
                      font_size=13, color_hex=self.WHITE, align="center")

    def add_vision_slide(self, vm):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "理念・ビジョン", "STEP 7")
        rows = [
            ("経営理念", vm.philosophy or ""),
            ("ビジョン",  vm.vision    or ""),
            ("ミッション", vm.mission  or ""),
            ("バリュー",  vm.values    or ""),
        ]
        y = 1.5
        for label, val in rows:
            self._rect(slide, 0.4, y, 2.2, 0.55, self.ACCENT)
            self._textbox(slide, 0.4, y, 2.2, 0.55, label,
                          font_size=13, bold=True, color_hex=self.WHITE, align="center")
            self._textbox(slide, 2.8, y, 10.0, 0.55, val,
                          font_size=13, color_hex=self.TEXT)
            y += 0.8

    def add_market_slide(self, ext):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "市場概況", "STEP 2")
        items = []
        if ext.market_size:
            items.append(f"市場規模: {ext.market_size}")
        if ext.market_growth:
            items.append(f"市場成長率: {ext.market_growth}")
        if ext.competitive_intensity:
            items.append(f"業界魅力度（5フォース）: {ext.competitive_intensity}")
        if ext.key_trends:
            items.append("主要トレンド・機会:")
            items += [f"  {t}" for t in ext.key_trends[:6]]
        self._bullets(slide, items)

    def add_pest_slide(self, ext):
        raw_pest = ext.raw.get("pest", {})
        if not raw_pest:
            return
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "PEST(EL)分析", "STEP 2-3")
        headers = ["カテゴリ", "事実・観察事項", "影響度", "機会/脅威"]
        rows = []
        cat_labels = {
            "political":     "政治",
            "economic":      "経済",
            "social":        "社会",
            "technological": "技術",
            "environmental": "環境",
            "legal":         "法規",
        }
        for key, label in cat_labels.items():
            for r in (raw_pest.get(key) or []):
                fact = r.get("事実・観察事項", "")
                if not fact:
                    continue
                rows.append([label, fact[:60], r.get("影響度", ""), r.get("機会/脅威", "")])
        if rows:
            self._table_slide(slide, headers, rows[:18])

    def add_five_forces_slide(self, ext):
        raw_ff = ext.raw.get("five_forces", {})
        if not raw_ff:
            return
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "5フォース分析", "STEP 2-3")
        force_labels = {
            "rivalry":          "競合他社の脅威",
            "new_entrants":     "新規参入の脅威",
            "substitutes":      "代替品の脅威",
            "buyer_power":      "買い手の交渉力",
            "supplier_power":   "売り手の交渉力",
        }
        items = []
        for key, label in force_labels.items():
            fd = raw_ff.get(key, {})
            if isinstance(fd, dict):
                score   = fd.get("score", "")
                summary = fd.get("summary", "")
                items.append(f"{label}: {score}/5　{summary}")
        items.append(f"総合業界魅力度: {raw_ff.get('overall_attractiveness', '')}")
        self._bullets(slide, items)

    def add_competitors_slide(self, ext):
        raw_comps = ext.raw.get("competitors", [])
        if not raw_comps:
            return
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "競合分析", "STEP 2-3")
        headers = ["企業名", "市場シェア", "主な強み", "脅威レベル"]
        rows = []
        for c in raw_comps[:12]:
            rows.append([
                c.get("name", ""),
                c.get("share", ""),
                c.get("strengths", "")[:50],
                c.get("threat_level", ""),
            ])
        if rows:
            self._table_slide(slide, headers, rows)

    def add_swot_slide(self, swot):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "SWOT分析", "STEP 8")
        # 4-quadrant layout
        quad = [
            ("S 強み",    swot.strengths,    0.3,  1.4, self.ACCENT),
            ("W 弱み",    swot.weaknesses,   6.9,  1.4, "#dc2626"),
            ("O 機会",    swot.opportunities, 0.3,  4.4, "#16a34a"),
            ("T 脅威",    swot.threats,      6.9,  4.4, "#d97706"),
        ]
        for label, items, lft, tp, color in quad:
            self._rect(slide, lft, tp, 6.2, 0.45, color)
            self._textbox(slide, lft, tp, 6.2, 0.45, label,
                          font_size=13, bold=True, color_hex=self.WHITE, align="center")
            text = "\n".join(f"・{x}" for x in (items or [])[:4])
            self._textbox(slide, lft, tp + 0.5, 6.2, 2.5, text,
                          font_size=11, color_hex=self.TEXT)

    def add_cross_swot_slide(self, swot):
        """クロスSWOT戦略スライド"""
        strategies = [
            ("SO戦略（攻め）", swot.so_strategies),
            ("ST戦略（差別化）", swot.st_strategies),
            ("WO戦略（強化）", swot.wo_strategies),
            ("WT戦略（防御）", swot.wt_strategies),
        ]
        has_any = any(s for _, s in strategies)
        if not has_any:
            return
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "クロスSWOT戦略", "STEP 8")
        colors = [self.ACCENT, "#dc2626", "#16a34a", "#d97706"]
        positions = [(0.3, 1.4), (6.9, 1.4), (0.3, 4.4), (6.9, 4.4)]
        for (label, items), color, (lft, tp) in zip(strategies, colors, positions):
            self._rect(slide, lft, tp, 6.2, 0.45, color)
            self._textbox(slide, lft, tp, 6.2, 0.45, label,
                          font_size=12, bold=True, color_hex=self.WHITE, align="center")
            text = "\n".join(f"・{x}" for x in (items or [])[:3])
            self._textbox(slide, lft, tp + 0.5, 6.2, 2.5, text,
                          font_size=11, color_hex=self.TEXT)

    def add_root_cause_slide(self, rc):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "真因分析", "STEP 9")
        items = []
        if rc.root_issue:
            items.append(f"【主課題】 {rc.root_issue}")
        if rc.primary_symptom:
            items.append(f"【主症状】 {rc.primary_symptom}")
        if rc.likely_root_causes:
            items.append("【根本原因】")
            items += [f"  {x}" for x in rc.likely_root_causes]
        self._bullets(slide, items)

    def add_strategy_slide(self, strategy):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "全社戦略仮説", "STEP 10")
        items = []
        if strategy.selected_strategy_name:
            items.append(f"【選定戦略】 {strategy.selected_strategy_name}")
        if strategy.selected_strategy_description:
            items.append(f"【内容】 {strategy.selected_strategy_description}")
        if strategy.selected_strategy_rationale:
            items.append(f"【根拠】 {strategy.selected_strategy_rationale}")
        if strategy.all_options:
            items.append(f"検討選択肢: {', '.join(strategy.all_options)}")
        self._bullets(slide, items)

    def add_domain_slide(self, domain):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "事業ドメイン設定", "STEP 11")
        rows = [
            ("ドメイン文",  domain.domain_statement or ""),
            ("ターゲット顧客", domain.customer or ""),
            ("提供価値",   domain.value_prop or ""),
            ("提供方法",   domain.method or ""),
            ("競争優位源泉", domain.competitive_source or ""),
        ]
        y = 1.5
        for label, val in rows:
            self._rect(slide, 0.4, y, 2.5, 0.5, self.NAVY)
            self._textbox(slide, 0.4, y, 2.5, 0.5, label,
                          font_size=12, bold=True, color_hex=self.WHITE, align="center")
            self._textbox(slide, 3.1, y, 9.8, 0.5, val,
                          font_size=12, color_hex=self.TEXT)
            y += 0.75

    def add_financial_slide(self, fin):
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._header_bar(slide, "財務サマリー", "STEP 4")
        headers = ["指標", "値"]
        rows = []
        if fin.revenue:
            rows.append(["売上高", f"{fin.revenue:,.0f} 百万円"])
        if fin.operating_profit:
            rows.append(["営業利益", f"{fin.operating_profit:,.0f} 百万円"])
        if fin.operating_margin is not None:
            rows.append(["営業利益率", f"{fin.operating_margin:.1f}%"])
        if fin.net_profit:
            rows.append(["当期純利益", f"{fin.net_profit:,.0f} 百万円"])
        if fin.total_assets:
            rows.append(["総資産", f"{fin.total_assets:,.0f} 百万円"])
        if fin.equity:
            rows.append(["純資産", f"{fin.equity:,.0f} 百万円"])
        if fin.roa is not None:
            rows.append(["ROA", f"{fin.roa:.1f}%"])
        if rows:
            self._table_slide(slide, headers, rows, top=1.5, font_size=14)

    def get_bytes(self) -> bytes:
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()


# ------------------------------------------------------------------ #
#  DOCX builder
# ------------------------------------------------------------------ #

def _build_docx(ctx: ProjectContext, company_name: str, sections: list) -> bytes:
    from core.docx_writer import DocxWriter
    w = DocxWriter()
    date_str = datetime.date.today().strftime("%Y年%m月%d日")
    w.add_title(f"中期経営計画 策定支援レポート\n{company_name}")
    w.document.add_paragraph(f"作成日: {date_str}")

    if "vision" in sections and ctx.vision_mission:
        vm = ctx.vision_mission
        w.add_chapter("1. 理念・ビジョン")
        for label, val in [("経営理念", vm.philosophy), ("ビジョン", vm.vision),
                            ("ミッション", vm.mission), ("バリュー", vm.values)]:
            if val:
                w.add_section(label)
                w.document.add_paragraph(val)

    if "market" in sections and ctx.external:
        ext = ctx.external
        w.add_chapter("2. 市場概況")
        items = []
        if ext.market_size:
            items.append(f"市場規模: {ext.market_size}")
        if ext.market_growth:
            items.append(f"市場成長率: {ext.market_growth}")
        if ext.competitive_intensity:
            items.append(f"業界魅力度: {ext.competitive_intensity}")
        w.add_content_from_markdown("\n".join(f"- {i}" for i in items))
        if ext.key_trends:
            w.add_section("主要トレンド")
            w.add_content_from_markdown("\n".join(f"- {t}" for t in ext.key_trends))

    if "pest" in sections and ctx.external:
        raw_pest = ctx.external.raw.get("pest", {})
        if raw_pest:
            w.add_chapter("3. PEST(EL)分析")
            cat_labels = {
                "political": "政治", "economic": "経済",
                "social": "社会", "technological": "技術",
                "environmental": "環境", "legal": "法規",
            }
            for key, label in cat_labels.items():
                rows = [r for r in (raw_pest.get(key) or []) if r.get("事実・観察事項")]
                if rows:
                    w.add_section(label)
                    for r in rows:
                        w.add_content_from_markdown(
                            f"- {r['事実・観察事項']}（影響度: {r.get('影響度','')} / {r.get('機会/脅威','')}）"
                        )

    if "five_forces" in sections and ctx.external:
        raw_ff = ctx.external.raw.get("five_forces", {})
        if raw_ff:
            w.add_chapter("4. 5フォース分析")
            force_labels = {
                "rivalry": "競合他社", "new_entrants": "新規参入",
                "substitutes": "代替品", "buyer_power": "買い手",
                "supplier_power": "売り手",
            }
            for key, label in force_labels.items():
                fd = raw_ff.get(key, {})
                if isinstance(fd, dict) and fd.get("summary"):
                    w.add_content_from_markdown(
                        f"- {label}: スコア{fd.get('score','')}/5 — {fd.get('summary','')}"
                    )
            w.document.add_paragraph(f"総合業界魅力度: {raw_ff.get('overall_attractiveness','')}")

    if "competitors" in sections and ctx.external:
        raw_comps = ctx.external.raw.get("competitors", [])
        if raw_comps:
            w.add_chapter("5. 競合分析")
            for c in raw_comps:
                w.add_section(c.get("name", ""))
                parts = []
                if c.get("share"):
                    parts.append(f"市場シェア: {c['share']}")
                if c.get("strengths"):
                    parts.append(f"強み: {c['strengths']}")
                if c.get("threat_level"):
                    parts.append(f"脅威レベル: {c['threat_level']}")
                w.add_content_from_markdown("\n".join(f"- {p}" for p in parts))

    if "swot" in sections and ctx.swot:
        swot = ctx.swot
        w.add_chapter("6. SWOT分析")
        for label, items in [("強み (S)", swot.strengths), ("弱み (W)", swot.weaknesses),
                               ("機会 (O)", swot.opportunities), ("脅威 (T)", swot.threats)]:
            if items:
                w.add_section(label)
                w.add_content_from_markdown("\n".join(f"- {x}" for x in items))
        if any([swot.so_strategies, swot.st_strategies, swot.wo_strategies, swot.wt_strategies]):
            w.add_section("クロスSWOT戦略")
            for label, items in [("SO戦略", swot.so_strategies), ("ST戦略", swot.st_strategies),
                                   ("WO戦略", swot.wo_strategies), ("WT戦略", swot.wt_strategies)]:
                if items:
                    w.add_content_from_markdown(f"**{label}**\n" + "\n".join(f"- {x}" for x in items))

    if "root_cause" in sections and ctx.root_cause:
        rc = ctx.root_cause
        w.add_chapter("7. 真因分析")
        if rc.root_issue:
            w.add_section("主課題")
            w.document.add_paragraph(rc.root_issue)
        if rc.primary_symptom:
            w.add_section("主症状")
            w.document.add_paragraph(rc.primary_symptom)
        if rc.likely_root_causes:
            w.add_section("根本原因")
            w.add_content_from_markdown("\n".join(f"- {x}" for x in rc.likely_root_causes))

    if "strategy" in sections and ctx.strategy:
        st = ctx.strategy
        w.add_chapter("8. 全社戦略仮説")
        if st.selected_strategy_name:
            w.add_section(st.selected_strategy_name)
        if st.selected_strategy_description:
            w.document.add_paragraph(st.selected_strategy_description)
        if st.selected_strategy_rationale:
            w.add_section("選定根拠")
            w.document.add_paragraph(st.selected_strategy_rationale)

    if "domain" in sections and ctx.domain:
        d = ctx.domain
        w.add_chapter("9. 事業ドメイン")
        for label, val in [
            ("ドメイン文", d.domain_statement), ("ターゲット顧客", d.customer),
            ("提供価値", d.value_prop), ("提供方法", d.method),
            ("競争優位源泉", d.competitive_source),
        ]:
            if val:
                w.add_content_from_markdown(f"**{label}**: {val}")

    if "financial" in sections and ctx.financial:
        fin = ctx.financial
        w.add_chapter("10. 財務サマリー")
        items = []
        if fin.revenue:
            items.append(f"売上高: {fin.revenue:,.0f} 百万円")
        if fin.operating_profit:
            items.append(f"営業利益: {fin.operating_profit:,.0f} 百万円")
        if fin.operating_margin is not None:
            items.append(f"営業利益率: {fin.operating_margin:.1f}%")
        if fin.net_profit:
            items.append(f"当期純利益: {fin.net_profit:,.0f} 百万円")
        if fin.roa is not None:
            items.append(f"ROA: {fin.roa:.1f}%")
        w.add_content_from_markdown("\n".join(f"- {i}" for i in items))

    buf = io.BytesIO()
    w.document.save(buf)
    return buf.getvalue()


# ------------------------------------------------------------------ #
#  PPTX builder
# ------------------------------------------------------------------ #

def _build_pptx(ctx: ProjectContext, company_name: str, sections: list) -> bytes:
    ppt = ConsultingPPT(company_name)
    date_str = datetime.date.today().strftime("%Y年%m月%d日")
    ppt.add_title_slide(date_str)

    if "vision" in sections and ctx.vision_mission:
        ppt.add_vision_slide(ctx.vision_mission)
    if "financial" in sections and ctx.financial:
        ppt.add_financial_slide(ctx.financial)
    if "market" in sections and ctx.external:
        ppt.add_market_slide(ctx.external)
    if "pest" in sections and ctx.external:
        ppt.add_pest_slide(ctx.external)
    if "five_forces" in sections and ctx.external:
        ppt.add_five_forces_slide(ctx.external)
    if "competitors" in sections and ctx.external:
        ppt.add_competitors_slide(ctx.external)
    if "swot" in sections and ctx.swot:
        ppt.add_swot_slide(ctx.swot)
        ppt.add_cross_swot_slide(ctx.swot)
    if "root_cause" in sections and ctx.root_cause:
        ppt.add_root_cause_slide(ctx.root_cause)
    if "strategy" in sections and ctx.strategy:
        ppt.add_strategy_slide(ctx.strategy)
    if "domain" in sections and ctx.domain:
        ppt.add_domain_slide(ctx.domain)

    return ppt.get_bytes()


# ------------------------------------------------------------------ #
#  Page UI
# ------------------------------------------------------------------ #

SECTION_DEFS = [
    ("vision",      "理念・ビジョン",    "vision_mission"),
    ("financial",   "財務サマリー",       "financial"),
    ("market",      "市場概況",           "external"),
    ("pest",        "PEST(EL)分析",       "external"),
    ("five_forces", "5フォース分析",      "external"),
    ("competitors", "競合分析",           "external"),
    ("swot",        "SWOT分析",           "swot"),
    ("root_cause",  "真因分析",           "root_cause"),
    ("strategy",    "全社戦略仮説",       "strategy"),
    ("domain",      "事業ドメイン",       "domain"),
]


def main():
    st.set_page_config(page_title="成果物エクスポート", page_icon="📤", layout="wide",
    initial_sidebar_state="expanded")
    if not check_auth():
        st.warning("ログインが必要です")
        return
    render_sidebar()

    client_id = st.session_state.get("client_id")
    if not client_id:
        st.warning("プロジェクトを選択してください。")
        if st.button("← プロジェクト一覧"):
            st.switch_page("app.py")
        return

    st.title("📤 成果物エクスポート")
    st.caption("蓄積された分析データを PowerPoint / Word に自動組み立てします")

    # ---- load context ----
    with st.spinner("データ読み込み中..."):
        ctx = ProjectContext.load(client_id)
    company_name = ctx.company_name or client_id[:8]

    # ---- availability summary ----
    summary = ctx.to_summary_dict()
    avail_map = {
        "vision_mission": summary["has_vision"],
        "financial":      summary["has_financial"],
        "external":       summary["has_external"],
        "swot":           summary["has_swot"],
        "root_cause":     summary["has_root_cause"],
        "strategy":       summary["has_strategy"],
        "domain":         summary["has_domain"],
    }

    st.subheader(f"📁 {company_name}")
    done = summary["completed_steps"]
    st.progress(done / 18, text=f"完了: {done}/18 ステップ")

    col_info, col_select = st.columns([1, 1])

    with col_info:
        st.markdown("**利用可能データ**")
        for key, label in [
            ("has_financial", "財務データ"), ("has_external", "外部環境"),
            ("has_vision",    "理念・ビジョン"), ("has_swot", "SWOT"),
            ("has_root_cause", "真因分析"), ("has_strategy", "戦略仮説"),
            ("has_domain",    "ドメイン"),
        ]:
            icon = "✅" if summary[key] else "○"
            st.markdown(f"{icon} {label}")

    with col_select:
        st.markdown("**含めるセクション**")
        selected_sections = []
        for sec_id, label, data_key in SECTION_DEFS:
            has_data = avail_map.get(data_key, False)
            disabled = not has_data
            default  = has_data
            checked  = st.checkbox(
                label + ("" if has_data else "  *(データなし)*"),
                value=default,
                disabled=disabled,
                key=f"sec_{sec_id}",
            )
            if checked:
                selected_sections.append(sec_id)

    st.divider()

    if not selected_sections:
        st.info("少なくとも1つのセクションを選択してください。")
        return

    col_ppt, col_doc = st.columns(2)

    with col_ppt:
        st.markdown("#### PowerPoint (pptx)")
        if st.button("🎨 PPTXを生成", use_container_width=True, type="primary"):
            with st.spinner("スライド生成中..."):
                try:
                    pptx_bytes = _build_pptx(ctx, company_name, selected_sections)
                    fname = f"{company_name}_report_{datetime.date.today():%Y%m%d}.pptx"
                    st.download_button(
                        label="⬇️ PPTXダウンロード",
                        data=pptx_bytes,
                        file_name=fname,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )
                    st.success(f"{len(selected_sections) + 1} スライドを生成しました")
                except Exception as e:
                    st.error(f"PPTX生成エラー: {e}")
                    import traceback
                    st.code(traceback.format_exc())

    with col_doc:
        st.markdown("#### Word (docx)")
        if st.button("📝 DOCXを生成", use_container_width=True):
            with st.spinner("文書生成中..."):
                try:
                    docx_bytes = _build_docx(ctx, company_name, selected_sections)
                    fname = f"{company_name}_report_{datetime.date.today():%Y%m%d}.docx"
                    st.download_button(
                        label="⬇️ DOCXダウンロード",
                        data=docx_bytes,
                        file_name=fname,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                    )
                    st.success("文書を生成しました")
                except Exception as e:
                    st.error(f"DOCX生成エラー: {e}")
                    import traceback
                    st.code(traceback.format_exc())


if __name__ == "__main__":
    main()
else:
    main()
